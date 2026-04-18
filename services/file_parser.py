# One function that detects the file type and routes to the right parser
# All parsers return plain text — the rest of the pipeline doesn't care about format

import requests
import tempfile
import os

def download_and_parse(pdf_url: str) -> str:
    # 1. Download the file
    response = requests.get(pdf_url)
    response.raise_for_status()

    # 2. Detect extension from the URL
    ext = os.path.splitext(pdf_url.split("?")[0])[-1].lower()

    # 3. Save to temp file
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as f:
        f.write(response.content)
        tmp_path = f.name

    # 4. Route to the right parser based on extension
    parsers = {
        ".pdf":  _parse_pdf,
        ".txt":  _parse_txt,
        ".md":   _parse_txt,   # markdown is just text
        ".docx": _parse_docx,
        ".doc":  _parse_docx,  # treated same as docx
        ".xlsx": _parse_xlsx,
        ".xls":  _parse_xlsx,
        ".pptx": _parse_pptx,
        ".ppt":  _parse_pptx,
        ".rtf":  _parse_rtf,
        ".csv":  _parse_txt,   # csv is plain text
    }

    parser = parsers.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file type: {ext}")

    return parser(tmp_path)


# ── Individual parsers ──────────────────────────────────────────────

def _parse_pdf(path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)


def _parse_txt(path: str) -> str:
    # Plain text, markdown, CSV — just read it
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    # Each paragraph is one line; join them all
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            # Convert each row to tab-separated text, skip fully empty rows
            row_text = "\t".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"[Slide {i+1}]")
        for shape in slide.shapes:
            # Only text-bearing shapes (titles, text boxes, tables)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _parse_rtf(path: str) -> str:
    from striprtf.striprtf import rtf_to_text
    with open(path, "r", errors="ignore") as f:
        return rtf_to_text(f.read())